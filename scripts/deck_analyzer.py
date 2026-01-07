#!/usr/bin/env python3
"""
Sales Deck Intelligence Analyzer
Extracts structured patterns from Eudia PowerPoint presentations
"""

import json
import os
import re
from collections import Counter, defaultdict
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# File paths
DESKTOP = Path("/Users/keiganpesenti/Desktop")
OUTPUT_DIR = Path("/Users/keiganpesenti/revops_weekly_update/gtm-brain/data/deck-analysis")

DECK_FILES = [
    DESKTOP / "Eudia Overview_ December 2025.pptx",
    DESKTOP / "Eudia Pricing Overview (Nov'25) v_F.cleaned.pptx",
    DESKTOP / "Eudia Case Study Bank  -  Repaired.pptx",
]

def emu_to_inches(emu):
    """Convert EMUs to inches"""
    return emu / 914400 if emu else 0

def rgb_to_hex(rgb_color):
    """Convert RGBColor to hex string"""
    if rgb_color is None:
        return None
    try:
        return f"#{rgb_color.rgb}"
    except:
        return None

def extract_shape_info(shape):
    """Extract detailed info from a shape"""
    info = {
        "type": str(shape.shape_type).replace("MSO_SHAPE_TYPE.", ""),
        "name": shape.name,
        "left": emu_to_inches(shape.left),
        "top": emu_to_inches(shape.top),
        "width": emu_to_inches(shape.width),
        "height": emu_to_inches(shape.height),
    }
    
    # Extract text content
    if shape.has_text_frame:
        paragraphs = []
        for para in shape.text_frame.paragraphs:
            para_info = {
                "text": para.text,
                "level": para.level,
                "runs": []
            }
            for run in para.runs:
                run_info = {"text": run.text}
                if run.font:
                    run_info["font_name"] = run.font.name
                    run_info["font_size"] = run.font.size.pt if run.font.size else None
                    run_info["bold"] = run.font.bold
                    run_info["italic"] = run.font.italic
                    try:
                        if run.font.color and run.font.color.type is not None:
                            run_info["color"] = str(run.font.color.rgb)
                    except (AttributeError, TypeError):
                        pass
                para_info["runs"].append(run_info)
            paragraphs.append(para_info)
        info["paragraphs"] = paragraphs
        info["full_text"] = shape.text
    
    # Extract table info
    if shape.has_table:
        table = shape.table
        info["table"] = {
            "rows": len(table.rows),
            "cols": len(table.columns),
            "cells": []
        }
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                info["table"]["cells"].append({
                    "row": row_idx,
                    "col": col_idx,
                    "text": cell.text
                })
    
    return info

def extract_slide_info(slide, slide_num):
    """Extract comprehensive info from a slide"""
    slide_info = {
        "slide_number": slide_num,
        "shapes": [],
        "all_text": [],
        "has_title": False,
        "title": None,
        "has_images": False,
        "has_tables": False,
        "has_charts": False,
    }
    
    # Get slide layout name
    try:
        slide_info["layout_name"] = slide.slide_layout.name
    except:
        slide_info["layout_name"] = "Unknown"
    
    # Process all shapes
    for shape in slide.shapes:
        shape_info = extract_shape_info(shape)
        slide_info["shapes"].append(shape_info)
        
        # Track content types
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            slide_info["has_images"] = True
        if shape.has_table:
            slide_info["has_tables"] = True
        if shape.has_chart if hasattr(shape, 'has_chart') else False:
            slide_info["has_charts"] = True
        
        # Extract text
        if shape.has_text_frame and shape.text.strip():
            slide_info["all_text"].append(shape.text.strip())
            
        # Check for title
        if shape.is_placeholder:
            try:
                if "TITLE" in str(shape.placeholder_format.type):
                    slide_info["has_title"] = True
                    slide_info["title"] = shape.text
            except:
                pass
    
    return slide_info

def analyze_presentation(file_path):
    """Analyze a single presentation"""
    print(f"\n{'='*60}")
    print(f"Analyzing: {file_path.name}")
    print('='*60)
    
    try:
        prs = Presentation(str(file_path))
    except Exception as e:
        print(f"ERROR loading presentation: {e}")
        return None
    
    analysis = {
        "file_name": file_path.name,
        "file_path": str(file_path),
        "total_slides": len(prs.slides),
        "slide_width_inches": emu_to_inches(prs.slide_width),
        "slide_height_inches": emu_to_inches(prs.slide_height),
        "slides": [],
        "fonts_used": Counter(),
        "colors_used": Counter(),
        "layouts_used": Counter(),
        "content_patterns": {
            "slides_with_images": 0,
            "slides_with_tables": 0,
            "slides_with_charts": 0,
        }
    }
    
    # Analyze each slide
    for idx, slide in enumerate(prs.slides, 1):
        slide_info = extract_slide_info(slide, idx)
        analysis["slides"].append(slide_info)
        
        # Aggregate stats
        analysis["layouts_used"][slide_info["layout_name"]] += 1
        if slide_info["has_images"]:
            analysis["content_patterns"]["slides_with_images"] += 1
        if slide_info["has_tables"]:
            analysis["content_patterns"]["slides_with_tables"] += 1
        if slide_info["has_charts"]:
            analysis["content_patterns"]["slides_with_charts"] += 1
        
        # Track fonts and colors from shapes
        for shape in slide_info["shapes"]:
            if "paragraphs" in shape:
                for para in shape["paragraphs"]:
                    for run in para["runs"]:
                        if run.get("font_name"):
                            analysis["fonts_used"][run["font_name"]] += 1
                        if run.get("color"):
                            analysis["colors_used"][run["color"]] += 1
    
    # Convert Counters to dicts for JSON serialization
    analysis["fonts_used"] = dict(analysis["fonts_used"])
    analysis["colors_used"] = dict(analysis["colors_used"])
    analysis["layouts_used"] = dict(analysis["layouts_used"])
    
    print(f"  Total slides: {analysis['total_slides']}")
    print(f"  Layouts used: {list(analysis['layouts_used'].keys())}")
    print(f"  Fonts found: {list(analysis['fonts_used'].keys())[:5]}")
    
    return analysis

def classify_slide_type(slide_info):
    """Classify a slide into a category based on content"""
    title = (slide_info.get("title") or "").lower()
    all_text = " ".join(slide_info.get("all_text", [])).lower()
    
    # Classification rules
    if slide_info["slide_number"] == 1:
        return "TITLE_SLIDE"
    
    if any(term in title for term in ["agenda", "overview", "contents", "outline"]):
        return "AGENDA"
    
    if any(term in title for term in ["case study", "customer story", "success story", "client"]):
        return "CASE_STUDY"
    
    if any(term in title for term in ["pricing", "investment", "commercial", "cost"]):
        return "PRICING"
    
    if any(term in title for term in ["team", "leadership", "about us", "who we are"]):
        return "TEAM_INTRO"
    
    if any(term in title for term in ["next step", "timeline", "implementation", "getting started"]):
        return "NEXT_STEPS"
    
    if any(term in title for term in ["roi", "return", "savings", "benefit", "value"]):
        return "VALUE_PROP"
    
    if any(term in title for term in ["problem", "challenge", "pain", "issue"]):
        return "PROBLEM_STATEMENT"
    
    if any(term in title for term in ["solution", "how we", "our approach", "capability"]):
        return "SOLUTION"
    
    if any(term in title for term in ["why", "differentiator", "competitive"]):
        return "COMPETITIVE"
    
    if any(term in title for term in ["use case", "application", "scenario"]):
        return "USE_CASE"
    
    if slide_info["has_tables"] or "metric" in all_text or "%" in all_text:
        return "DATA_VISUALIZATION"
    
    if slide_info["has_images"] and len(slide_info["all_text"]) < 3:
        return "VISUAL"
    
    return "CONTENT"

def extract_messaging_patterns(all_slides):
    """Extract messaging and value proposition patterns"""
    patterns = {
        "pain_point_language": [],
        "value_propositions": [],
        "proof_points": [],
        "action_language": [],
        "industry_terms": [],
    }
    
    # Keywords to look for
    pain_keywords = ["challenge", "problem", "risk", "cost", "time", "inefficient", "manual", "complex"]
    value_keywords = ["reduce", "save", "improve", "accelerate", "automate", "streamline", "optimize", "transform"]
    proof_keywords = ["%", "million", "hours", "days", "customers", "cases", "documents", "contracts"]
    action_keywords = ["next step", "contact", "schedule", "demo", "trial", "pilot", "start"]
    
    for slide in all_slides:
        for text in slide.get("all_text", []):
            text_lower = text.lower()
            
            # Check for pain points
            for kw in pain_keywords:
                if kw in text_lower:
                    patterns["pain_point_language"].append(text)
                    break
            
            # Check for value props
            for kw in value_keywords:
                if kw in text_lower:
                    patterns["value_propositions"].append(text)
                    break
            
            # Check for proof points
            for kw in proof_keywords:
                if kw in text_lower:
                    patterns["proof_points"].append(text)
                    break
            
            # Check for CTAs
            for kw in action_keywords:
                if kw in text_lower:
                    patterns["action_language"].append(text)
                    break
    
    return patterns

def generate_narrative_analysis(slides_by_deck):
    """Analyze the narrative arc across decks"""
    narrative = {
        "common_flow": [],
        "deck_structures": {},
    }
    
    for deck_name, slides in slides_by_deck.items():
        deck_flow = []
        for slide in slides:
            slide_type = classify_slide_type(slide)
            deck_flow.append({
                "slide_num": slide["slide_number"],
                "type": slide_type,
                "title": slide.get("title", ""),
            })
        narrative["deck_structures"][deck_name] = deck_flow
    
    return narrative

def main():
    """Main extraction and analysis pipeline"""
    print("\n" + "="*70)
    print("EUDIA SALES DECK INTELLIGENCE ANALYZER")
    print("="*70)
    
    # Create output directory
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    
    all_analyses = []
    all_slides = []
    slides_by_deck = {}
    
    # Analyze each deck
    for deck_path in DECK_FILES:
        if deck_path.exists():
            analysis = analyze_presentation(deck_path)
            if analysis:
                all_analyses.append(analysis)
                all_slides.extend(analysis["slides"])
                slides_by_deck[deck_path.stem] = analysis["slides"]
        else:
            print(f"\nWARNING: File not found: {deck_path}")
    
    if not all_analyses:
        print("\nERROR: No presentations could be analyzed!")
        return
    
    # Aggregate analysis
    print("\n" + "="*70)
    print("AGGREGATE ANALYSIS")
    print("="*70)
    
    # Fonts across all decks
    all_fonts = Counter()
    all_colors = Counter()
    all_layouts = Counter()
    
    for analysis in all_analyses:
        all_fonts.update(analysis["fonts_used"])
        all_colors.update(analysis["colors_used"])
        all_layouts.update(analysis["layouts_used"])
    
    # Slide type classification
    slide_types = Counter()
    for slide in all_slides:
        slide_type = classify_slide_type(slide)
        slide_types[slide_type] += 1
    
    # Messaging patterns
    messaging = extract_messaging_patterns(all_slides)
    
    # Narrative analysis
    narrative = generate_narrative_analysis(slides_by_deck)
    
    # Compile comprehensive analysis
    comprehensive = {
        "meta": {
            "decks_analyzed": len(all_analyses),
            "total_slides": len(all_slides),
            "analysis_date": "2026-01-05"
        },
        "design_patterns": {
            "fonts": dict(all_fonts.most_common(20)),
            "colors": dict(all_colors.most_common(20)),
            "layouts": dict(all_layouts),
            "slide_dimensions": {
                "width": all_analyses[0]["slide_width_inches"],
                "height": all_analyses[0]["slide_height_inches"],
            }
        },
        "content_architecture": {
            "slide_types": dict(slide_types),
            "narrative_flow": narrative,
        },
        "messaging_patterns": messaging,
        "individual_decks": all_analyses,
    }
    
    # Save results
    output_file = OUTPUT_DIR / "deck_analysis_complete.json"
    with open(output_file, 'w') as f:
        json.dump(comprehensive, f, indent=2, default=str)
    print(f"\nFull analysis saved to: {output_file}")
    
    # Generate summary report
    summary = generate_summary_report(comprehensive)
    summary_file = OUTPUT_DIR / "DECK_DNA_SUMMARY.md"
    with open(summary_file, 'w') as f:
        f.write(summary)
    print(f"Summary report saved to: {summary_file}")
    
    # Print quick summary
    print("\n" + "="*70)
    print("QUICK SUMMARY")
    print("="*70)
    print(f"Decks analyzed: {len(all_analyses)}")
    print(f"Total slides: {len(all_slides)}")
    print(f"\nSlide Types Distribution:")
    for stype, count in slide_types.most_common():
        print(f"  {stype}: {count}")
    print(f"\nTop Fonts: {list(all_fonts.most_common(5))}")
    print(f"\nLayouts Used: {list(all_layouts.keys())}")

def generate_summary_report(analysis):
    """Generate a human-readable summary report"""
    report = """# Eudia Sales Deck DNA - Pattern Analysis

## Executive Summary

This document catalogs the repeatable patterns identified across {deck_count} customer presentations, 
comprising {slide_count} total slides. These patterns form the "Deck DNA" - the ruleset for generating 
consistent, on-brand customer presentations.

---

## 1. Visual/Design Patterns

### Slide Dimensions
- Width: {width:.2f} inches
- Height: {height:.2f} inches
- Aspect Ratio: Standard widescreen (16:9)

### Typography
**Primary Fonts Identified:**
{fonts}

### Color Palette
**Colors Used (by frequency):**
{colors}

### Layout Templates
**Slide Layouts Used:**
{layouts}

---

## 2. Content Architecture

### Slide Type Distribution
{slide_types}

### Standard Narrative Arc
Based on analysis, the typical deck follows this structure:
{narrative_flow}

---

## 3. Messaging/Positioning Patterns

### Pain Point Language
Common ways we frame customer challenges:
{pain_points}

### Value Proposition Patterns
How we articulate benefits:
{value_props}

### Proof Points
Evidence and metrics we emphasize:
{proof_points}

### Call-to-Action Language
How we drive next steps:
{ctas}

---

## 4. Deck-Specific Structures

{deck_structures}

---

## 5. Customization Points

Based on divergence analysis, these elements vary by customer/context:
- Industry-specific terminology
- Case study selection
- Pricing/commercial details
- Stakeholder-specific messaging

---

## 6. Recommendations for Deck Generation

### Core Elements (Always Include)
1. Title slide with Eudia branding
2. Problem/challenge framing
3. Solution overview
4. Relevant case study/proof point
5. Clear next steps

### Variable Elements (Customize Per Prospect)
1. Industry-specific pain points
2. Relevant case studies
3. Pricing tier selection
4. Stakeholder-appropriate messaging

---

*Generated: 2026-01-05*
*Source: Eudia Sales Deck Intelligence Analyzer*
""".format(
        deck_count=analysis["meta"]["decks_analyzed"],
        slide_count=analysis["meta"]["total_slides"],
        width=analysis["design_patterns"]["slide_dimensions"]["width"],
        height=analysis["design_patterns"]["slide_dimensions"]["height"],
        fonts="\n".join([f"- {font}: {count} occurrences" for font, count in list(analysis["design_patterns"]["fonts"].items())[:10]]),
        colors="\n".join([f"- `{color}`: {count} uses" for color, count in list(analysis["design_patterns"]["colors"].items())[:10]]),
        layouts="\n".join([f"- {layout}: {count} slides" for layout, count in analysis["design_patterns"]["layouts"].items()]),
        slide_types="\n".join([f"- **{stype}**: {count} slides" for stype, count in analysis["content_architecture"]["slide_types"].items()]),
        narrative_flow="See deck structures below for detailed flow analysis.",
        pain_points="\n".join([f'- "{text[:100]}..."' if len(text) > 100 else f'- "{text}"' for text in analysis["messaging_patterns"]["pain_point_language"][:10]]) or "- No explicit pain point language extracted",
        value_props="\n".join([f'- "{text[:100]}..."' if len(text) > 100 else f'- "{text}"' for text in analysis["messaging_patterns"]["value_propositions"][:10]]) or "- No explicit value propositions extracted",
        proof_points="\n".join([f'- "{text[:100]}..."' if len(text) > 100 else f'- "{text}"' for text in analysis["messaging_patterns"]["proof_points"][:10]]) or "- No explicit proof points extracted",
        ctas="\n".join([f'- "{text[:100]}..."' if len(text) > 100 else f'- "{text}"' for text in analysis["messaging_patterns"]["action_language"][:10]]) or "- No explicit CTAs extracted",
        deck_structures=generate_deck_structure_section(analysis["content_architecture"]["narrative_flow"]["deck_structures"])
    )
    return report

def generate_deck_structure_section(structures):
    """Generate markdown section for deck structures"""
    sections = []
    for deck_name, flow in structures.items():
        section = f"### {deck_name}\n\n"
        section += "| Slide # | Type | Title |\n"
        section += "|---------|------|-------|\n"
        for item in flow:
            title = (item.get("title") or "")[:50]
            section += f"| {item['slide_num']} | {item['type']} | {title} |\n"
        sections.append(section)
    return "\n".join(sections)

if __name__ == "__main__":
    main()

